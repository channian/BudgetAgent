"""
Pipeline Step 1 — 案件初步分析與分類（規則 + 本地 LLM，零未知）
══════════════════════════════════════════════════════════════
流程：先用關鍵字規則判定系統；規則無法判定的，交給本地 Ollama LLM 判定。
      LLM 保證回傳白名單系統（重試 → 模糊比對 → 兜底），因此輸出 100% 已分配，
      不會再有「未知系統」，pensieve 完全退出此步驟。

輸入：網路資料夾（每個子資料夾 = 一個案件）
輸出：
  - pytollm_output.json   (全部案件，皆已分配系統)
  - system_defined.json   (= 全部案件)
  - system_unknown.json   (恆為空陣列，僅保留檔案相容性)

支援檔案：.pdf / .pptx / .docx / .txt / .xlsx / .xlsm
          （Excel 整列以「 | 」串接，金額與項目綁同一行，便於 LLM 判斷預算佔比）

前置：
  • Ollama 已啟動，且模型已下載（ollama pull qwen2.5:7b）
  • 套件：pip install openpyxl python-pptx PyPDF2 python-docx

下一步：執行 pipeline_2.py（本地 LLM 初步審核並直接寫入 DB）。
"""

import os, json, re

# ── 路徑設定 ────────────────────────────────────────────────────────────
INPUT_PATH = r"\\khfs4\4600$\01.Leon Yi\01.FAC2\預算\Leon WBS專案預算簽核系統\預算審核"
WORK_DIR   = r"D:\ASEKH\K20076\2026\預算AI Agent\新思路0409\初步審核"

# ── 系統關鍵字對照表（依「廠務設備邏輯決策官」白名單 v2）──────────────
# 白名單：水務 / 空壓 / 空調 / 電力 / 安全 / 消防 / 資安 / AI自動化 /
#         抽氣 / Relayout / 二次配 / 建管
# 註：監控已併入「資安」（不再是獨立系統）；二次配無固定關鍵字 — 由衝突
#     解決規則或 LLM 判定（見下方規則）。
ENTITY_MAPPING = {
    "水務":    ["MBR", "廢水", "回收水", "RO", "超純水", "純水", "生活污水", "污水泵", "沉水馬達", "SiO2分析儀", "矽分析儀", "藥槽區", "污泥", "加藥", "水箱", "熱泵(節能水)", "給排水"],
    "空壓":    ["空壓機", "乾燥機", "CDA", "真空", "外熱式乾燥機"],
    "空調":    ["冷卻水塔", "冰水機", "PCW", "MAU", "FFU", "溫濕度", "空調備援", "冷氣", "熱泵(空調系)", "換氣風機"],
    "電力":    ["電盤", "變壓器", "UPS", "發電機", "GCB", "斷路器", "電力主系統", "配電箱", "照明", "變電站"],
    "安全":    ["公危倉", "化學品泄露", "緊急應變", "洗眼器", "防護衣"],
    "消防":    ["防火門", "消防栓", "偵煙器", "廣播系統", "灑水頭", "消防法規", "消防泵"],
    "資安":    ["SCADA", "PLC", "自動化監測", "中控系統", "CCTV", "MOF電錶", "資訊安全", "資安", "網路安全", "防火牆", "弱點掃描", "滲透測試", "入侵偵測", "EDR", "SIEM", "ISO27001", "存取控制", "資安監控", "SOC"],
    "AI自動化":["影像辨識模型", "深度學習", "AI演算法"],
    "抽氣":    ["RRTO", "RTO", "沸石滾輪", "Scrubber", "洗滌塔", "有機排氣", "無機排氣", "排氣管路", "風車", "異味改善", "抽氣馬達"],
    "建管":    ["建置案", "餐廳建置", "老舊建物改善", "園管區", "結構安全", "防水工程"],
    "Relayout":["機台移位", "空間規劃", "隔間拆除", "裝修", "家具配置", "佈局調整"],
}

# ── 衝突解決優先序觸發詞（高 → 低，對應 classification_prompt.md）────────
RELAYOUT_TRIGGERS = ["家具配置", "家具", "隔間", "空間規劃", "空間調整", "空間",
                     "裝修", "拆除", "拆機", "機台進駐", "機台移位", "機台放置",
                     "CLASS增設", "RELAYOUT", "佈局調整", "佈局"]
ERJIPEI_TRIGGER   = "二次配"
JIANGUAN_TRIGGERS = ["法規", "結構安全", "違章"]
GENERIC_PARTS     = ["法蘭", "蝶閥", "管帽"]

# ── 本地 LLM 設定（Ollama）──────────────────────────────────────────────
# pipeline_1 = 規則判定 + 本地 LLM 補完，輸出保證「零未知」，不再經 pensieve。
#   主力推薦：qwen2.5:7b（夠準）；想更穩 → qwen2.5:14b；3B 級太小不建議當主力。
LLM_MODEL      = "qwen2.5:7b"           # Ollama 已下載的模型名稱（ollama pull <name>）
OLLAMA_URL     = "http://localhost:11434"
LLM_MAX_CHARS  = 3000                    # 每案件最多傳給 LLM 的文字長度（節省 token）
LLM_RETRIES    = 2                       # LLM 回答無效時的重試次數
# 兜底：LLM 連線/解析全失敗時的最終分類（保證不留未知）。
# 依 classification_prompt.md「通用/模糊一律歸 Relayout」精神，預設 Relayout。
LLM_FALLBACK_SYSTEM = "Relayout"

LLM_WHITELIST = ["水務", "空壓", "空調", "電力", "安全", "消防", "資安",
                 "AI自動化", "抽氣", "Relayout", "二次配", "建管"]

EXPERT_DB = {
    "設備擴充 (UTI)": {"空調": ["黃金燦"], "空壓": ["郭于斌"], "水務": ["梁益齊"], "抽氣": ["梁益齊"], "電力": ["李明鴻"], "資安": ["王嘉漢"]},
    "工程擴廠 (新工)": {"Relayout": ["陳信舟", "鄭仁勝", "紀志忠"], "二次配": ["陳信舟", "鄭仁勝", "紀志忠"], "空調": ["鄭仁勝"], "水務": ["陳妍方"], "抽氣": ["陳妍方"]},
    "CIM相關":  {"資安": ["王嘉漢"], "AI自動化": ["黃互慶"]},
    "法遵 (ESH)": {"消防": ["吳明華"], "建管": ["吳明華"], "環保": ["姜婷毓"]},
}

BOILERPLATE_NOISE = [
    "Need To Know", "會議機密", "四「不」原則", "不將會議記錄轉寄",
    "公務機密管理規範", "人事規章處理", "雇用合約", "Notes E-Mail",
    "公司之財產", "資安關鍵字", "非我自身公務不過問", "不擅自邀請",
    "不揭露產品功能", "不傳遞信件", "禁止作為私人用途",
]


def _load_system_prompt() -> str:
    """Read classification_prompt.md from same directory."""
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "classification_prompt.md")
    with open(path, encoding="utf-8") as f:
        return f.read()


# 金額訊號：貨幣符號/關鍵字、千分位數字、帶單位數字、5+位純數字（排除 4 位年份）
_MONEY_RE = re.compile(
    r"NT\$|US\$|新台幣|新臺幣|金額|報價|預算|總價|單價|小計|合計|總計|費用|含稅|未稅"
    r"|\d{1,3}(?:,\d{3})+(?:\.\d+)?"
    r"|\d+(?:\.\d+)?\s*(?:萬元|萬|億|元)"
    r"|\d{5,}"
)


def _extract_amount_digest(text: str, max_lines: int = 40, max_line_len: int = 120) -> str:
    """
    從完整全文掃出「含金額」的行，集中成一個 digest（去重、限長）。
    讓 LLM 判斷各系統預算佔比（rule 7）時，金額一目了然、與項目綁在同一行。
    """
    seen, lines = set(), []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line or line.startswith("<<<") or line.startswith("[工作表"):
            continue
        if _MONEY_RE.search(line):
            snippet = line[:max_line_len]
            if snippet not in seen:
                seen.add(snippet)
                lines.append(f"- {snippet}")
                if len(lines) >= max_lines:
                    break
    return "\n".join(lines)


def _llm_call(case: dict, system_prompt: str, strict: bool = False) -> str:
    """Send one case to Ollama, return the raw model text ('' on connection error)."""
    import urllib.request
    import urllib.error

    full_text       = case.get("所有檔案原文(按檔案分段)", "")
    content_snippet = full_text[:LLM_MAX_CHARS]
    # 金額線索從「完整全文」擷取，避免被 3000 字截斷漏掉後段金額
    money_digest    = _extract_amount_digest(full_text)

    user_msg = (
        f"案件名稱：{case['案件名稱']}\n"
        f"檔案清單：\n{case.get('檔案實體清單', '')}\n"
    )
    if money_digest:
        user_msg += (
            "\n【金額線索】（自動擷取，供判斷各系統預算佔比/ROI 參考；"
            "若多系統並存，請依金額最高者判定）\n"
            f"{money_digest}\n"
        )
    user_msg += f"\n文件內容（節錄）：\n{content_snippet}"

    if strict:
        user_msg += (
            "\n\n【重要】只能輸出以下其中一個詞，不得有任何其他文字："
            + "、".join(LLM_WHITELIST)
        )

    payload = json.dumps({
        "model": LLM_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_msg},
        ],
        "stream": False,
        "options": {"temperature": 0},
    }, ensure_ascii=False).encode("utf-8")

    req = urllib.request.Request(
        f"{OLLAMA_URL}/api/chat",
        data=payload,
        headers={"Content-Type": "application/json"},
    )
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            result = json.loads(resp.read())
        return (result.get("message", {}).get("content") or "").strip()
    except urllib.error.URLError:
        print(f"    ⚠️  無法連線 Ollama（{OLLAMA_URL}）— 請確認 Ollama 已啟動且模型已下載")
        return ""
    except Exception as e:
        print(f"    ⚠️  LLM 呼叫失敗：{e}")
        return ""


def _normalize_to_whitelist(raw: str) -> str:
    """Map raw LLM text to a valid whitelist category, or '' if none recognised."""
    if not raw:
        return ""
    cleaned = raw.strip("【】「」《》()（）.。,，、:：\n\r\t ")
    tokens = cleaned.split()
    if tokens and tokens[0] in LLM_WHITELIST:
        return tokens[0]
    if cleaned in LLM_WHITELIST:
        return cleaned
    # 模糊比對：白名單詞出現在回應任何位置
    for cat in LLM_WHITELIST:
        if cat in raw:
            return cat
    return ""


def _llm_classify(case: dict, system_prompt: str) -> str:
    """
    保證回傳一個白名單系統：重試 → 模糊比對 → 兜底 fallback，絕不回「未知」。
    """
    for attempt in range(LLM_RETRIES):
        raw = _llm_call(case, system_prompt, strict=(attempt > 0))
        sysname = _normalize_to_whitelist(raw)
        if sysname:
            return sysname
        if raw:
            print(f"    ↻ 第 {attempt + 1} 次回應無效「{raw[:30]}」，重試…")
    print(f"    ↪ 連續無效，套用兜底分類：{LLM_FALLBACK_SYSTEM}")
    return LLM_FALLBACK_SYSTEM


def extract_text(file_path):
    ext = os.path.splitext(file_path)[-1].lower()
    try:
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            return "\n".join(
                shape.text for slide in prs.slides
                for shape in slide.shapes if hasattr(shape, "text")
            )
        if ext == ".pdf":
            from PyPDF2 import PdfReader
            return "".join(p.extract_text() or "" for p in PdfReader(file_path).pages)
        if ext == ".docx":
            from docx import Document
            return "\n".join(p.text for p in Document(file_path).paragraphs)
        if ext == ".txt":
            with open(file_path, encoding="utf-8") as f:
                return f.read()
        if ext in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            out = []
            for ws in wb.worksheets:
                out.append(f"[工作表: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    # 整列用「 | 」串起來，讓「項目 ↔ 金額」綁在同一行，便於金額判斷
                    cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                    if cells:
                        out.append(" | ".join(cells))
            wb.close()
            return "\n".join(out)
    except Exception:
        pass
    return "[注意：此檔案內容提取受限]"


def _systems_in_text(text_up):
    """回傳 {系統: 關鍵字最後出現位置}，供「名稱含多系統取後者」判斷使用。"""
    found = {}
    for sys, kws in ENTITY_MAPPING.items():
        pos = max((text_up.rfind(kw.upper()) for kw in kws), default=-1)
        if pos >= 0:
            found[sys] = pos
    return found


def find_best_system(content, folder_name):
    """依「廠務設備邏輯決策官」衝突解決優先序判斷唯一系統；無法判定回傳「未知」。"""
    raw_name = folder_name or ""
    raw_all  = f"{folder_name} {content}"
    name_up  = raw_name.upper()
    text_up  = raw_all.upper()
    for noise in BOILERPLATE_NOISE:
        text_up = text_up.replace(noise.upper(), "")

    # 1) Relayout 最高優先：空間/裝修/家具/拆除/機台進駐 …
    if any(t.upper() in text_up for t in RELAYOUT_TRIGGERS):
        return "Relayout"

    # 2) 二次配 次高優先
    if ERJIPEI_TRIGGER in raw_all:
        return "二次配"

    # 3) 建管：僅限法規 / 結構安全 / 違章
    if any(t in raw_all for t in JIANGUAN_TRIGGERS):
        return "建管"

    # 6) 專案名稱含二個系統 → 以後者（出現位置較後）為準
    name_systems = _systems_in_text(name_up)
    if len(name_systems) >= 2:
        return max(name_systems, key=name_systems.get)

    # 關鍵字頻率計分
    scores  = {s: sum(text_up.count(kw.upper()) for kw in kws)
               for s, kws in ENTITY_MAPPING.items()}
    nonzero = {s: sc for s, sc in scores.items() if sc > 0}

    # 4) 通用組件（法蘭/蝶閥/管帽）且無明確主體 → Relayout
    if not nonzero and any(g in raw_all for g in GENERIC_PARTS):
        return "Relayout"

    if not nonzero:
        # 名稱命中單一系統也採用
        if len(name_systems) == 1:
            return next(iter(name_systems))
        return "未知"

    # 5) 多系統且含 Relayout → Relayout
    if len(nonzero) >= 2 and "Relayout" in nonzero:
        return "Relayout"

    ranked = sorted(nonzero.items(), key=lambda x: x[1], reverse=True)
    best, best_score = ranked[0]
    second_score     = ranked[1][1] if len(ranked) > 1 else 0

    # 7) 多系統分數相同無法分辨 → 交給 LLM（依各系統預算金額最高者判定）
    if second_score == best_score:
        return "未知"
    return best


def classify_category(system, folder_and_header):
    # CIM 相關：資安（含原監控）/ AI自動化
    if system in ("AI自動化", "資安"):
        return "CIM相關"
    # 法遵 (ESH)：消防 / 建管 / 安全
    if system in ("消防", "建管", "安全"):
        return "法遵 (ESH)"
    # 工程擴廠 (新工)：二次配 / Relayout 一律歸新工
    if system in ("二次配", "Relayout"):
        return "工程擴廠 (新工)"
    if system == "未知":
        return "未知"
    # 其餘設備系統：依關鍵字判斷新建工程或既有設備擴充
    new_keywords = ["修繕", "建置", "新建", "整改", "備援", "RELAYOUT", "擴建", "新增"]
    return "工程擴廠 (新工)" if any(k in folder_and_header for k in new_keywords) else "設備擴充 (UTI)"


def process():
    os.makedirs(WORK_DIR, exist_ok=True)
    folders = [d for d in os.listdir(INPUT_PATH)
               if os.path.isdir(os.path.join(INPUT_PATH, d)) and not d.startswith(".")]

    all_cases = []
    for folder in folders:
        print(f"  分析中：{folder}")
        case_path = os.path.join(INPUT_PATH, folder)
        files     = [f for f in os.listdir(case_path) if not f.startswith("~")]

        inventory = []
        full_text = ""
        for fname in files:
            inventory.append(f"- {fname}")
            ext = os.path.splitext(fname)[1].lower()
            if ext in (".pdf", ".pptx", ".docx", ".txt", ".xlsx", ".xlsm"):
                extracted = extract_text(os.path.join(case_path, fname))
                full_text += f"\n<<< 檔案名稱: {fname} 開始 >>>\n{extracted}\n<<< 檔案名稱: {fname} 結束 >>>\n"
            else:
                full_text += f"\n<<< 檔案名稱: {fname} (非文字檔) >>>\n"

        system   = find_best_system(full_text, folder)
        header   = (folder + full_text[:1500]).upper()
        category = classify_category(system, header)
        experts  = EXPERT_DB.get(category, {}).get(system, [])

        all_cases.append({
            "案件名稱":            folder,
            "判定系統":            system,
            "判定類別":            category,
            "負責專家":            ", ".join(experts) if experts else "待分配",
            "檔案實體清單":         "\n".join(inventory),
            "所有檔案原文(按檔案分段)": full_text,
        })

    # ── 本地 LLM 補完規則無法判定的案件（取代 pensieve）────────────────
    rule_unknown = [c for c in all_cases if c["判定系統"] == "未知"]
    if rule_unknown:
        print(f"\n🤖 規則判定剩 {len(rule_unknown)} 筆未知 → 本地 LLM（{LLM_MODEL}）補完…")
        try:
            system_prompt = _load_system_prompt()
        except FileNotFoundError:
            print("❌ 找不到 classification_prompt.md，將全部套用兜底分類")
            system_prompt = None

        for case in rule_unknown:
            print(f"  分類中：{case['案件名稱']}")
            if system_prompt:
                sys_result = _llm_classify(case, system_prompt)
            else:
                sys_result = LLM_FALLBACK_SYSTEM
            # _llm_classify 保證回傳白名單系統，絕不留未知
            case["判定系統"] = sys_result
            folder_header = (
                case["案件名稱"] + case.get("所有檔案原文(按檔案分段)", "")[:1500]
            ).upper()
            case["判定類別"] = classify_category(sys_result, folder_header)
            experts = EXPERT_DB.get(case["判定類別"], {}).get(sys_result, [])
            case["負責專家"] = ", ".join(experts) if experts else "待分配"
            print(f"    → ✅ {sys_result} / {case['判定類別']}")

    # 此時 all_cases 已 100% 分配完成，無未知系統
    rule_count = len(all_cases) - len(rule_unknown)
    _save(os.path.join(WORK_DIR, "pytollm_output.json"), all_cases)
    _save(os.path.join(WORK_DIR, "system_defined.json"), all_cases)
    _save(os.path.join(WORK_DIR, "system_unknown.json"), [])   # 恆為空，保留檔案相容性

    print(f"\n✅ 完成：共 {len(all_cases)} 筆全部判定完成"
          f"（規則 {rule_count} 筆 / LLM {len(rule_unknown)} 筆），無未知系統")
    print("   → 直接執行 pipeline_2.py（本地審核並寫入 DB）")


def _save(path, data):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)


if __name__ == "__main__":
    process()
